# -*- coding: utf-8 -*-
"""
Windows 친화 MCP Agent: 문서 내용 검색기
- 툴: search_docs
- 대상 확장자: .pdf .docx .pptx .xlsx .txt .md
- 필요시 .docx/.pptx → PDF로 변환하여 (docx2pdf) Downloads/<원본>/CONVERTED 에 저장 후 검색
- Excel(.xlsx)은 openpyxl로 직접 내용 검색 (변환 불필요)
- PDF 텍스트는 PyMuPDF로 추출

인자 예시:
  call_tool("search_docs", {
    "query": "인공지능",
    "root_dir": null,                # null이면 기본: 사용자 Downloads
    "convert_to_pdf": true,          # .docx/.pptx 변환 시도
    "case_sensitive": false,
    "use_regex": false,
    "max_results": 50,
    "include_exts": [".pdf",".docx",".pptx",".xlsx",".txt",".md"]
  })
"""

import os
import re
import sys
import json
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple

from mcp.server.fastmcp import FastMCP

mcp = FastMCP("doc-search-agent")

# ---------------------------
# 공통 유틸
# ---------------------------
DEFAULT_EXTS = [".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"]

def get_downloads_dir() -> Path:
    return Path.home() / "Downloads"

def ensure_converted_dir(src_file: Path) -> Path:
    out_dir = src_file.parent / "CONVERTED"
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir

def norm_ext(p: Path) -> str:
    return p.suffix.lower()

def read_small_text_file(p: Path, encoding_list=("utf-8", "cp949", "euc-kr")) -> str:
    for enc in encoding_list:
        try:
            return p.read_text(encoding=enc, errors="ignore")
        except Exception:
            continue
    # 마지막 시도: 바이너리 열고 디코드
    try:
        return p.read_bytes().decode("utf-8", errors="ignore")
    except Exception:
        return ""

# ---------------------------
# 변환(.docx/.pptx → .pdf) - 선택적
# ---------------------------
def try_convert_office_to_pdf(src: Path) -> Optional[Path]:
    """
    docx/pptx를 docx2pdf로 PDF 변환 시도(Word/PowerPoint 설치 필요)
    성공 시 PDF 경로 반환, 실패 시 None
    """
    ext = norm_ext(src)
    if ext not in {".docx", ".pptx"}:
        return None

    try:
        from docx2pdf import convert  # 내부에서 COM 사용 (Windows + 오피스 설치 필요)
    except Exception:
        return None

    out_dir = ensure_converted_dir(src)
    out_pdf = out_dir / (src.stem + ".pdf")
    if out_pdf.exists():
        return out_pdf

    # docx2pdf는 "파일 하나" 또는 "디렉토리" 단위로 동작하므로, 임시로 파일만 변환
    # 일부 환경에서 from / to 인자 지원이 제한적일 수 있어 원본 디렉토리에 생성 후 이동하는 전략 사용
    try:
        # 원본 위치에서 변환 수행 (docx2pdf는 결과를 같은 폴더에 생성)
        convert(str(src))
        created = src.with_suffix(".pdf")
        if created.exists():
            created.replace(out_pdf)
            return out_pdf
        # 일부 버전에선 바로 같은 폴더에 생성되지 않을 수 있음 → 폴더 단위 변환 시도
        convert(str(src.parent))
        if created.exists():
            created.replace(out_pdf)
            return out_pdf
    except Exception:
        return None

    return None

# ---------------------------
# 텍스트 추출기
# ---------------------------
def extract_text_pdf(p: Path) -> List[Tuple[int, str]]:
    """
    PDF에서 페이지별 텍스트를 반환 [(page_no, text), ...] (1-based page)
    """
    try:
        import fitz  # PyMuPDF
    except Exception:
        return []

    pages = []
    try:
        with fitz.open(str(p)) as doc:
            for i, page in enumerate(doc, start=1):
                text = page.get_text("text") or ""
                pages.append((i, text))
    except Exception:
        pass
    return pages

def extract_text_docx(p: Path) -> str:
    try:
        import docx
    except Exception:
        return ""
    try:
        d = docx.Document(str(p))
        texts = []
        for para in d.paragraphs:
            texts.append(para.text)
        # 표 내용도 간단히 포함
        for table in d.tables:
            for row in table.rows:
                cells = [c.text for c in row.cells]
                texts.append("\t".join(cells))
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text_pptx(p: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""
    try:
        prs = Presentation(str(p))
        texts = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text_xlsx(p: Path, max_sheet_rows: int = 1000) -> str:
    """
    Excel을 텍스트로 평탄화 (상위 일정 행까지만 – 너무 큰 파일 안전장치)
    """
    try:
        import openpyxl
    except Exception:
        return ""
    try:
        wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
        texts = []
        for ws in wb.worksheets:
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                row_count += 1
                if row_count > max_sheet_rows:
                    break
                cells = [("" if v is None else str(v)) for v in row]
                line = "\t".join(cells).strip()
                if line:
                    texts.append(f"[{ws.title}] {line}")
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text_generic(p: Path) -> str:
    ext = norm_ext(p)
    if ext == ".pdf":
        # PDF는 페이지별로 처리 → 여기서는 전체 합쳐서 반환
        pages = extract_text_pdf(p)
        return "\n".join([t for _, t in pages])
    elif ext == ".docx":
        return extract_text_docx(p)
    elif ext == ".pptx":
        return extract_text_pptx(p)
    elif ext == ".xlsx":
        return extract_text_xlsx(p)
    elif ext in {".txt", ".md"}:
        return read_small_text_file(p)
    else:
        return ""

def extract_snippets(text: str, pattern: str, use_regex: bool, case_sensitive: bool,
                     context: int = 60, max_hits: int = 5) -> List[str]:
    flags = 0 if case_sensitive else re.IGNORECASE
    matches = []
    if use_regex:
        try:
            for m in re.finditer(pattern, text, flags):
                start = max(0, m.start() - context)
                end = min(len(text), m.end() + context)
                snippet = text[start:end].replace("\n", " ")
                matches.append(snippet)
                if len(matches) >= max_hits:
                    break
        except re.error:
            return []
    else:
        needle = pattern if case_sensitive else pattern.lower()
        hay = text if case_sensitive else text.lower()
        idx = 0
        while True:
            idx = hay.find(needle, idx)
            if idx == -1:
                break
            start = max(0, idx - context)
            end = min(len(text), idx + len(pattern) + context)
            snippet = text[start:end].replace("\n", " ")
            matches.append(snippet)
            if len(matches) >= max_hits:
                break
            idx += len(pattern)
    return matches

# ---------------------------
# 검색 핵심
# ---------------------------
def search_file(p: Path, query: str, use_regex: bool, case_sensitive: bool) -> Dict[str, Any]:
    ext = norm_ext(p)

    # PDF는 페이지별 스니펫 추출
    if ext == ".pdf":
        page_hits = []
        pages = extract_text_pdf(p)
        for page_no, page_text in pages:
            snippets = extract_snippets(page_text, query, use_regex, case_sensitive)
            if snippets:
                page_hits.append({"page": page_no, "snippets": snippets})
        return {"path": str(p), "ext": ext, "hits": page_hits}

    # 그 외는 전체 텍스트에서 스니펫
    full = extract_text_generic(p)
    snippets = extract_snippets(full, query, use_regex, case_sensitive)
    return {"path": str(p), "ext": ext, "hits": [{"page": None, "snippets": snippets}] if snippets else []}

# ---------------------------
# MCP Tool
# ---------------------------
@mcp.tool()
def search_docs(
    query: str,
    root_dir: Optional[str] = None,
    include_exts: Optional[List[str]] = None,
    convert_to_pdf: bool = True,
    case_sensitive: bool = False,
    use_regex: bool = False,
    max_results: int = 50
) -> dict:
    """
    문서 내용 검색 (필요 시 .docx/.pptx → PDF 변환 후 검색)
    - 기본 검색 경로: 사용자 Downloads
    - include_exts: 검색 대상 확장자 목록 (기본: .pdf .docx .pptx .xlsx .txt .md)
    - convert_to_pdf=True일 때 docx/pptx는 PDF 변환 시도(성공 시 PDF에서 검색)
    """
    try:
        if not query:
            return {"ok": False, "message": "query가 비어 있습니다.", "results": []}

        base = Path(root_dir).expanduser().resolve() if root_dir else get_downloads_dir()
        if not base.exists():
            return {"ok": False, "message": f"검색 루트가 없습니다: {base}", "results": []}

        exts = [e.lower() for e in (include_exts or DEFAULT_EXTS)]
        results: List[Dict[str, Any]] = []
        converted: List[str] = []
        scanned_count = 0

        for p in base.rglob("*"):
            if not p.is_file():
                continue
            ext = norm_ext(p)
            if ext not in exts:
                continue

            target_path = p
            # 변환: docx/pptx → PDF (가능할 때만)
            if convert_to_pdf and ext in {".docx", ".pptx"}:
                pdf_path = try_convert_office_to_pdf(p)
                if pdf_path and pdf_path.exists():
                    target_path = pdf_path
                    converted.append(str(pdf_path))

            # 검색
            hit = search_file(target_path, query, use_regex, case_sensitive)
            if hit.get("hits"):
                results.append(hit)

            scanned_count += 1
            if len(results) >= max_results:
                break

        summary = {
            "searched_root": str(base),
            "scanned_files": scanned_count,
            "converted_outputs": converted,
            "max_results": max_results
        }
        return {"ok": True, "message": "검색 완료", "summary": summary, "results": results}

    except Exception as e:
        return {"ok": False, "message": f"예외 발생: {e}", "results": []}

if __name__ == "__main__":
    mcp.run()
